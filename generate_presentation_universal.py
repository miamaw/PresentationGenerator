"""
Universal PowerPoint Generator
================================
Generic, customizable presentation generator for educators
No branding - fully customizable fonts, colors, and backgrounds
"""

import sys
import os
import json
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


class FlexibleLayoutEngine:
    """
    Intelligently stacks multiple layout sections on a single slide.
    
    Supports:
    - Content (single column text)
    - Left/Right (two columns)
    - LeftTop/RightTop/LeftBottom/RightBottom (four boxes)
    - Any combination of the above!
    """
    
    def __init__(self, slide_data, dimensions):
        """
        Args:
            slide_data: Dict with keys: content, left, right, left_top, etc.
            dimensions: Dict with CONTENT_LEFT, CONTENT_TOP, CONTENT_WIDTH, CONTENT_HEIGHT, etc.
        """
        self.data = slide_data
        self.dims = dimensions
        self.sections = []  # List of (type, priority, height_weight)
        
    def analyze_sections(self):
        """
        Detect which sections are present and determine layout strategy.
        
        Returns list of sections in order: [(section_type, data), ...]
        """
        sections = []
        
        # Check for Content section (header/intro text)
        if self.data.get("content"):
            sections.append({
                'type': 'content',
                'data': self.data["content"],
                'priority': 1,  # Content usually goes first
                'height_weight': 0.3  # Takes 30% of space by default
            })
        
        # Check for four-box section
        has_four_box = any([
            self.data.get("left_top"),
            self.data.get("right_top"),
            self.data.get("left_bottom"),
            self.data.get("right_bottom")
        ])
        
        if has_four_box:
            sections.append({
                'type': 'four_box',
                'data': {
                    'left_top': self.data.get("left_top", []),
                    'right_top': self.data.get("right_top", []),
                    'left_bottom': self.data.get("left_bottom", []),
                    'right_bottom': self.data.get("right_bottom", [])
                },
                'priority': 2,  # Main content section
                'height_weight': 0.7  # Takes 70% of space
            })
        
        # Check for two-column section
        elif self.data.get("left") or self.data.get("right"):
            sections.append({
                'type': 'two_column',
                'data': {
                    'left': self.data.get("left", []),
                    'right': self.data.get("right", [])
                },
                'priority': 2,
                'height_weight': 0.7
            })
        
        # Special case: Reading comprehension layout
        # (LeftTop = passage, LeftBottom = questions, no RightTop/RightBottom)
        if (self.data.get("left_top") and self.data.get("left_bottom") and 
            not self.data.get("right_top") and not self.data.get("right_bottom")):
            sections = [{
                'type': 'reading',
                'data': {
                    'passage': self.data["left_top"],
                    'questions': self.data["left_bottom"]
                },
                'priority': 1,
                'height_weight': 1.0  # Takes full height
            }]
        
        return sections
    
    def calculate_section_positions(self, sections, total_height):
        """
        Calculate top position and height for each section.
        
        Args:
            sections: List of section dicts
            total_height: Total available height (Inches)
            
        Returns:
            List of (section, top, height) tuples
        """
        if not sections:
            return []
        
        # Special case: Single section gets full height
        if len(sections) == 1:
            return [(sections[0], self.dims['CONTENT_TOP'], total_height)]
        
        # Multiple sections: Distribute height based on weights
        total_weight = sum(s['height_weight'] for s in sections)
        gap = Inches(0.3)  # Gap between sections
        total_gap = gap * (len(sections) - 1)
        usable_height = total_height - total_gap
        
        positions = []
        current_top = self.dims['CONTENT_TOP']
        
        for section in sections:
            # Calculate this section's height
            weight_ratio = section['height_weight'] / total_weight
            section_height = usable_height * weight_ratio
            
            positions.append((section, current_top, section_height))
            current_top += section_height + gap
        
        return positions
    
    def render_content_section(self, slide, top, height, data, config, add_textbox_func):
        """Render a single-column content section"""
        font_size = 22 if height > Inches(2) else 18
        add_textbox_func(
            slide, 
            self.dims['CONTENT_LEFT'], 
            top,
            self.dims['CONTENT_WIDTH'], 
            height, 
            data,
            font_size=font_size,
            label="Content",
            config=config
        )
    
    def render_two_column_section(self, slide, top, height, data, config, add_textbox_func):
        """Render a two-column section"""
        font_size = 20 if height > Inches(2) else 16
        
        add_textbox_func(
            slide,
            self.dims['CONTENT_LEFT'],
            top,
            self.dims['COLUMN_WIDTH'],
            height,
            data['left'],
            font_size=font_size,
            label="Left",
            config=config
        )
        
        add_textbox_func(
            slide,
            self.dims['CONTENT_LEFT'] + self.dims['COLUMN_WIDTH'] + self.dims['COLUMN_GAP'],
            top,
            self.dims['COLUMN_WIDTH'],
            height,
            data['right'],
            font_size=font_size,
            label="Right",
            config=config
        )
    
    def render_four_box_section(self, slide, top, height, data, config, add_textbox_func):
        """Render a four-box section"""
        half_height = (height - self.dims['ROW_GAP']) / 2
        font_size = 18 if height > Inches(3) else 14
        
        # Top row
        add_textbox_func(
            slide,
            self.dims['CONTENT_LEFT'],
            top,
            self.dims['COLUMN_WIDTH'],
            half_height,
            data['left_top'],
            font_size=font_size,
            label="LeftTop",
            config=config
        )
        
        add_textbox_func(
            slide,
            self.dims['CONTENT_LEFT'] + self.dims['COLUMN_WIDTH'] + self.dims['COLUMN_GAP'],
            top,
            self.dims['COLUMN_WIDTH'],
            half_height,
            data['right_top'],
            font_size=font_size,
            label="RightTop",
            config=config
        )
        
        # Bottom row
        add_textbox_func(
            slide,
            self.dims['CONTENT_LEFT'],
            top + half_height + self.dims['ROW_GAP'],
            self.dims['COLUMN_WIDTH'],
            half_height,
            data['left_bottom'],
            font_size=font_size,
            label="LeftBottom",
            config=config
        )
        
        add_textbox_func(
            slide,
            self.dims['CONTENT_LEFT'] + self.dims['COLUMN_WIDTH'] + self.dims['COLUMN_GAP'],
            top + half_height + self.dims['ROW_GAP'],
            self.dims['COLUMN_WIDTH'],
            half_height,
            data['right_bottom'],
            font_size=font_size,
            label="RightBottom",
            config=config
        )
    
    def render_reading_section(self, slide, top, height, data, config, add_textbox_func):
        """Render a reading comprehension section"""
        passage_height = height * 0.65
        gap = Inches(0.3)
        questions_height = height * 0.35 - gap
        
        add_textbox_func(
            slide,
            self.dims['CONTENT_LEFT'],
            top,
            self.dims['CONTENT_WIDTH'],
            passage_height,
            data['passage'],
            font_size=18,
            label="ReadingPassage",
            config=config
        )
        
        add_textbox_func(
            slide,
            self.dims['CONTENT_LEFT'],
            top + passage_height + gap,
            self.dims['CONTENT_WIDTH'],
            questions_height,
            data['questions'],
            font_size=16,
            label="ReadingQuestions",
            config=config
        )
    
    def render(self, slide, config, add_textbox_func):
        """
        Main rendering function.
        
        Args:
            slide: python-pptx slide object
            config: Configuration dict
            add_textbox_func: Function to add textboxes
        """
        sections = self.analyze_sections()
        
        if not sections:
            # Empty slide - just title
            return
        
        # Calculate positions
        positions = self.calculate_section_positions(
            sections, 
            self.dims['CONTENT_HEIGHT']
        )
        
        # Render each section
        for section, top, height in positions:
            section_type = section['type']
            section_data = section['data']
            
            if section_type == 'content':
                self.render_content_section(
                    slide, top, height, section_data, config, add_textbox_func
                )
            elif section_type == 'two_column':
                self.render_two_column_section(
                    slide, top, height, section_data, config, add_textbox_func
                )
            elif section_type == 'four_box':
                self.render_four_box_section(
                    slide, top, height, section_data, config, add_textbox_func
                )
            elif section_type == 'reading':
                self.render_reading_section(
                    slide, top, height, section_data, config, add_textbox_func
                )


# === DEFAULT CONFIG ===
DEFAULT_CONFIG = {
    "background_image": None,  # Can be None for solid color
    "background_color": [255, 255, 255],  # White
    "title_color": [0, 0, 0],  # Black
    "text_color": [64, 64, 64],  # Dark gray
    "font_name": "Arial",
    "title_font_name": "Arial",
    "slide_width": 13.33,
    "slide_height": 7.5,
    "enable_slide_numbers": True,
    "enable_overflow_warnings": True,
    "styles": {
        "vocabulary": {"font_size": 24, "color": [0, 128, 0], "bold": True},
        "question": {"font_size": 20, "color": [128, 0, 128], "bold": False},
        "answer": {"font_size": 18, "color": [128, 128, 128], "italic": True},
        "emphasis": {"font_size": 22, "color": [192, 0, 0], "bold": True}
    }
}


def merge_config(user_config, defaults=DEFAULT_CONFIG):
    """Merge user config with defaults"""
    config = defaults.copy()
    if user_config:
        config.update(user_config)
        # Merge nested styles
        if "styles" in user_config:
            config["styles"] = {**defaults["styles"], **user_config["styles"]}
    return config


# === OVERFLOW DETECTION ===
def check_text_overflow(text, font_size, width_inches, height_inches):
    """Estimates if text will overflow"""
    if hasattr(width_inches, 'inches'):
        width_inches = width_inches.inches
    if hasattr(height_inches, 'inches'):
        height_inches = height_inches.inches
    
    chars_per_inch = 72 / font_size * 2.5
    chars_per_line = int(width_inches * chars_per_inch)
    
    words = text.split()
    current_line_length = 0
    lines_needed = 1
    
    for word in words:
        word_length = len(word) + 1
        if current_line_length + word_length > chars_per_line:
            lines_needed += 1
            current_line_length = word_length
        else:
            current_line_length += word_length
    
    line_height = font_size / 72 * 1.3
    lines_available = int(height_inches / line_height)
    
    will_overflow = lines_needed > lines_available
    return will_overflow, lines_needed, lines_available


# === LIST DETECTION ===
def is_list_content(lines):
    """Detect if content should be formatted as a list"""
    if not lines:
        return False
    
    bullet_patterns = [r'^\s*[â€¢\-\*]', r'^\s*\d+\.', r'^\s*[a-z]\)', r'^\s*[A-Z]\.']
    matching = sum(1 for line in lines if any(re.match(p, line) for p in bullet_patterns))
    return matching >= len(lines) * 0.5


def clean_bullet_marker(text):
    """Remove common bullet markers from text"""
    text = re.sub(r'^\s*[â€¢\-\*]\s*', '', text)
    text = re.sub(r'^\s*\d+\.\s*', '', text)
    text = re.sub(r'^\s*[a-z]\)\s*', '', text)
    return text


# === QUESTION SPLITTING ===
def split_questions(text):
    """Robustly split multiple questions"""
    questions = re.split(r'\?\s*(?=\d+\.|\b[A-Z])', text)
    result = []
    
    for q in questions:
        q = q.strip()
        if q:
            if not q.endswith('?'):
                q += '?'
            result.append(q)
    
    return result if result else [text]


# === STYLE APPLICATION ===
def apply_style(paragraph, style_name, config):
    """Apply predefined style to a paragraph"""
    styles = config.get("styles", {})
    
    if style_name in styles:
        style = styles[style_name]
        paragraph.font.size = Pt(style.get("font_size", 22))
        
        if "color" in style:
            color = style["color"]
            paragraph.font.color.rgb = RGBColor(*color)
        
        paragraph.font.bold = style.get("bold", False)
        paragraph.font.italic = style.get("italic", False)


def parse_styled_text(text):
    """Parse text with inline style markers"""
    match = re.match(r'^\[(\w+)\]\s*(.+)', text)
    if match:
        return match.group(1), match.group(2)
    return None, text


# === MATH/SPECIAL CHARACTERS ===
def process_math(text):
    """Convert simple math notation to Unicode symbols"""
    superscripts = {'0': 'â°', '1': 'Â¹', '2': 'Â²', '3': 'Â³', '4': 'â´', 
                    '5': 'âµ', '6': 'â¶', '7': 'â·', '8': 'â¸', '9': 'â¹'}
    text = re.sub(r'\^(\d)', lambda m: superscripts.get(m.group(1), m.group(1)), text)
    
    subscripts = {'0': 'â‚€', '1': 'â‚', '2': 'â‚‚', '3': 'â‚ƒ', '4': 'â‚„',
                  '5': 'â‚…', '6': 'â‚†', '7': 'â‚‡', '8': 'â‚ˆ', '9': 'â‚‰'}
    text = re.sub(r'_(\d)', lambda m: subscripts.get(m.group(1), m.group(1)), text)
    
    replacements = {
        '<=': 'â‰¤', '>=': 'â‰¥', '!=': 'â‰ ', '~=': 'â‰ˆ',
        'alpha': 'Î±', 'beta': 'Î²', 'gamma': 'Î³', 'delta': 'Î´',
        'pi': 'Ï€', 'theta': 'Î¸', 'sigma': 'Ïƒ'
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    
    return text


# === VALIDATION ===
def validate_slide(slide_data, slide_num, config):
    """Validate slide data"""
    issues = []
    
    if not slide_data["title"]:
        issues.append(f"Slide {slide_num}: Missing title")
    elif len(slide_data["title"]) > 100:
        issues.append(f"Slide {slide_num}: Title very long ({len(slide_data['title'])} chars)")
    
    has_content = any([
        slide_data["content"], slide_data["left"], slide_data["right"],
        slide_data["left_top"], slide_data["right_top"],
        slide_data["left_bottom"], slide_data["right_bottom"]
    ])
    
    if not has_content:
        issues.append(f"Slide {slide_num}: No content defined")
    
    return issues


# === ADD TEXTBOX ===
def add_textbox(slide, left, top, width, height, lines, font_size=22, label=None, 
                config=None, v_align=MSO_ANCHOR.TOP):
    """Enhanced textbox with overflow detection, list formatting, and styling"""
    if not lines:
        return None
    
    if config is None:
        config = DEFAULT_CONFIG
    
    joined = " ".join(lines)
    joined = process_math(joined)
    text_length = len(joined)
    
    # Overflow detection
    if config.get("enable_overflow_warnings", True):
        try:
            w = width.inches if hasattr(width, 'inches') else width
            h = height.inches if hasattr(height, 'inches') else height
            overflow, needed, available = check_text_overflow(joined, font_size, w, h)
            if overflow:
                print(f"âš ï¸  Potential overflow in '{label}': needs {needed} lines, has {available}")
        except Exception:
            pass
    
    # Auto font-size for long text
    if text_length > 300:
        font_size = min(font_size, 18)
    if text_length > 500:
        font_size = min(font_size, 16)
    if text_length > 700:
        font_size = min(font_size, 14)
    if text_length > 1000:
        font_size = min(font_size, 12)
    
    # Handle [step] animations
    step_lines = [l for l in lines if "[step]" in l.lower()]
    if step_lines:
        return add_step_textboxes(slide, left, top, width, lines, font_size, label, config)
    
    # Detect list formatting
    is_list = is_list_content(lines)
    
    # Create textbox
    box = slide.shapes.add_textbox(left, top, width, height)
    if label:
        box.name = label
    
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = v_align
    
    # Add content
    first = True
    for item in lines:
        if not item.strip():
            continue
        
        style, text = parse_styled_text(item)
        text = process_math(text)
        
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        
        if is_list:
            text = clean_bullet_marker(text)
            p.level = 0
            p.text = text
        else:
            p.text = text
        
        p.font.name = config["font_name"]
        if style:
            apply_style(p, style, config)
        else:
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*config["text_color"])
        
        if text_length > 300:
            p.space_after = Pt(3)
        else:
            p.space_after = Pt(6)
    
    return box


def add_step_textboxes(slide, left, top, width, lines, font_size, label, config):
    """Create separate textboxes for each [step] line"""
    top_offset = top
    boxes = []
    
    for i, item in enumerate(lines):
        if not item.strip():
            continue
        
        text = re.sub(r'\[step\]\s*', '', item, flags=re.IGNORECASE)
        text = process_math(text)
        style, text = parse_styled_text(text)
        
        box = slide.shapes.add_textbox(left, top_offset, width, Inches(0.6))
        if label:
            box.name = f"{label}_Step{i+1}"
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = config["font_name"]
        
        if style:
            apply_style(p, style, config)
        else:
            p.font.size = Pt(font_size)
            p.font.color.rgb = RGBColor(*config["text_color"])
        
        boxes.append(box)
        top_offset += Inches(0.65)
    
    return boxes


# === ADD SLIDE NUMBER ===
def add_slide_number(slide, slide_num, total_slides, config):
    """Add slide number footer"""
    footer_text = f"{slide_num} / {total_slides}"
    
    footer = slide.shapes.add_textbox(
        Inches(config["slide_width"] - 1.5),
        Inches(config["slide_height"] - 0.5),
        Inches(1),
        Inches(0.3)
    )
    
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT


# === PARSER ===
def parse_content_file(filename):
    """Parse content file"""
    slides = []
    current = {
        "title": "", "content": [], "notes": [],
        "left": [], "right": [],
        "left_top": [], "right_top": [],
        "left_bottom": [], "right_bottom": [],
        "template": None
    }
    section = None

    with open(filename, "r", encoding="utf-8") as f:
        for line in f:
            line = line.rstrip()
            
            if not line or line == "---":
                continue

            if line.startswith("Slide "):
                if current["title"]:
                    slides.append(current)
                    current = {
                        "title": "", "content": [], "notes": [],
                        "left": [], "right": [],
                        "left_top": [], "right_top": [],
                        "left_bottom": [], "right_bottom": [],
                        "template": None
                    }
                section = None
                continue

            if line.startswith("Template:"):
                current["template"] = line.replace("Template:", "").strip()
                continue

            if line.startswith("Title:"):
                current["title"] = line.replace("Title:", "").strip()
                section = None
                continue
            elif line.startswith("Content:"):
                section = "content"
                text = line.replace("Content:", "").strip()
            elif line.startswith("Left:"):
                section = "left"
                text = line.replace("Left:", "").strip()
            elif line.startswith("Right:"):
                section = "right"
                text = line.replace("Right:", "").strip()
            elif line.startswith("LeftTop:"):
                section = "left_top"
                text = line.replace("LeftTop:", "").strip()
            elif line.startswith("RightTop:"):
                section = "right_top"
                text = line.replace("RightTop:", "").strip()
            elif line.startswith("LeftBottom:"):
                section = "left_bottom"
                text = line.replace("LeftBottom:", "").strip()
                if any(q in text for q in ["1.", "2.", "3."]):
                    questions = split_questions(text)
                    current["left_bottom"].extend(questions)
                    continue
            elif line.startswith("RightBottom:"):
                section = "right_bottom"
                text = line.replace("RightBottom:", "").strip()
            elif line.startswith("Notes:"):
                section = "notes"
                text = line.replace("Notes:", "").strip()
            else:
                text = line

            if section in current and text:
                current[section].append(text)

        if current["title"]:
            slides.append(current)

    return slides


# === BUILD PRESENTATION ===
def build_presentation(slides, output_name, config=None):
    """Build presentation with custom styling"""
    if config is None:
        config = DEFAULT_CONFIG
    
    prs = Presentation()
    prs.slide_width = Inches(config["slide_width"])
    prs.slide_height = Inches(config["slide_height"])
    
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height
    TITLE_LEFT = Inches(1.5)
    TITLE_TOP = Inches(0.6)
    CONTENT_LEFT = Inches(1.5)
    CONTENT_TOP = Inches(1.5)
    CONTENT_WIDTH = Inches(10.5)
    CONTENT_HEIGHT = Inches(5.0)
    COLUMN_GAP = Inches(0.4)
    ROW_GAP = Inches(0.3)
    COLUMN_WIDTH = (CONTENT_WIDTH - COLUMN_GAP) / 2
    
    total_slides = len(slides)
    
    for idx, s in enumerate(slides, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Background
        if config.get("background_image") and os.path.exists(config["background_image"]):
            slide.shapes.add_picture(config["background_image"], 0, 0,
                                    width=SLIDE_WIDTH, height=SLIDE_HEIGHT)
        else:
            # Solid color background
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*config.get("background_color", [255, 255, 255]))
        
        # Title
        title_box = slide.shapes.add_textbox(TITLE_LEFT, TITLE_TOP, CONTENT_WIDTH, Inches(0.8))
        tf = title_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = s["title"]
        p.font.name = config.get("title_font_name", config["font_name"])
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(*config["title_color"])
        
        # Flexible Layout Logic - Handles ANY combination
        dimensions = {
            'CONTENT_LEFT': CONTENT_LEFT,
            'CONTENT_TOP': CONTENT_TOP,
            'CONTENT_WIDTH': CONTENT_WIDTH,
            'CONTENT_HEIGHT': CONTENT_HEIGHT,
            'COLUMN_WIDTH': COLUMN_WIDTH,
            'COLUMN_GAP': COLUMN_GAP,
            'ROW_GAP': ROW_GAP
        }
        
        layout_engine = FlexibleLayoutEngine(s, dimensions)
        layout_engine.render(slide, config, add_textbox)
        
        # Add slide numbers
        if config.get("enable_slide_numbers", True):
            add_slide_number(slide, idx, total_slides, config)
        
        # Notes
        if s["notes"]:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            for note in s["notes"]:
                notes_tf.add_paragraph().text = f"â€¢ {note}"
    
    prs.save(output_name)
    print(f"âœ… Presentation created: {output_name}")


# === MAIN ===
if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python generate_presentation_universal.py <content.txt> [config.json]")
        sys.exit(1)
    
    input_file = sys.argv[1]
    config_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    config = DEFAULT_CONFIG
    if config_file and os.path.exists(config_file):
        with open(config_file, 'r') as f:
            user_config = json.load(f)
            config = merge_config(user_config)
    
    slides = parse_content_file(input_file)
    
    base_name = input_file.replace(".txt", "")
    output_name = base_name + "_presentation.pptx"
    
    build_presentation(slides, output_name, config)